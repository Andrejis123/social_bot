"""Layout drawing functions.

Each `draw_*` adds a fresh slide to the presentation and paints shapes directly.
No slide master, no named placeholders — see project_reports_spec.md for the
architecture decision (flexed 2026-05-26).

Dark-theme slides (Cover, Intro, Overview, CategorySection) use navy bg + white
text + rainbow footer.
Light-theme slides (Summary, AdditionalData) use white bg + navy text + a
brand-color strip under the heading.
"""
from __future__ import annotations

import io
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import theme
from .brand import Brand

# ────────────────────────────────────────────────────────────────────────
# Data structures consumed by the renderer
# ────────────────────────────────────────────────────────────────────────

@dataclass
class CategoryItem:
    """One distinct campaign/event/theme inside a category."""
    title: str           # short item title (used internally; not always rendered)
    narrative: str       # one-line caption shown below the image
    image_path: Path     # representative post image
    post_ids: list[str] = field(default_factory=list)
                         # all post.ids in the underlying cluster; the renderer
                         # uses these to build per-page narrative when the
                         # category paginates across multiple slides


@dataclass
class CategoryPreview:
    """One entry on the per-account Intro slide — name + one image from that category."""
    name: str
    image_path: Path


@dataclass
class MetricCard:
    """One cell on Summary / AdditionalData."""
    label: str           # subheading
    value: str           # large primary text (e.g. "12.3K")
    caption: str = ""    # optional small detail line
    value_aux: str = ""  # smaller, non-bold trailing detail rendered inline
                         # after the value (e.g. "(10 reels)" after "10 posts")


# ────────────────────────────────────────────────────────────────────────
# Helpers
# ────────────────────────────────────────────────────────────────────────

LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF)


def _clean_text(text: str) -> str:
    """Strip em-dashes and dash-separators from rendered text.

    See feedback_no_emdashes.md. Backstop only, source content should be
    authored without dash separators. Substitute: comma + space (least
    destructive). Compound words (cross-brand, mid-May) keep their internal
    hyphens; en-dashes in date ranges (25 April – 25 May) are preserved
    as conventional typography.
    """
    return (
        text.replace(" — ", ", ")
            .replace(" -- ", ", ")
            .replace(" - ", ", ")
            .replace("—", ",")
    )


def _split_caption(text: str, max_first_words: int = 4) -> str:
    """Break a caption into a short first line + rest, so per-image captions
    don't read as one continuous strip across the bottom of a category slide.

    Prefer splitting on the first natural pause (comma/period/colon/semicolon)
    within the first `max_first_words` words. Falls back to a hard split after
    3 words.
    """
    words = text.split()
    if len(words) <= 3:
        return text
    for i, w in enumerate(words[:max_first_words]):
        if w.endswith((",", ".", ":", ";")):
            return " ".join(words[: i + 1]) + "\n" + " ".join(words[i + 1 :])
    return " ".join(words[:3]) + "\n" + " ".join(words[3:])


def format_metric(value: int) -> str:
    """Abbreviate a metric count for circle display."""
    if value < 1000:
        return str(value)
    if value < 1_000_000:
        v = value / 1000.0
        return f"{v:.1f}K" if v % 1 else f"{int(v)}K"
    v = value / 1_000_000.0
    return f"{v:.1f}M" if v % 1 else f"{int(v)}M"


def _add_slide(prs: Presentation, bg_color: RGBColor):
    """Blank slide with the given solid background fill."""
    layout = prs.slide_layouts[6]  # 'Blank' in default Office theme
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    return slide


def _place_image(slide, img_path: Path, left, top, width, height,
                 *, fit: bool = False) -> None:
    """Center-crop image to box aspect, then place. Avoids the recreation's
    'whitespace borders inside cropped images' defect.

    When `fit=True`, instead of cropping, scale the source down to fit fully
    within the box and center it (letterbox). Use for logos / wordmarks where
    cropping would chop off text.
    """
    target_w_emu = int(width)
    target_h_emu = int(height)
    target_ratio = target_w_emu / target_h_emu

    with Image.open(img_path) as im:
        src_w, src_h = im.size
        src_ratio = src_w / src_h

        if fit:
            # Letterbox: shrink to fit, center inside the box.
            if src_ratio > target_ratio:
                placed_w = width
                placed_h = int(width / src_ratio)
            else:
                placed_h = height
                placed_w = int(height * src_ratio)
            placed_left = left + (target_w_emu - int(placed_w)) // 2
            placed_top = top + (target_h_emu - int(placed_h)) // 2
            slide.shapes.add_picture(
                str(img_path), placed_left, placed_top,
                width=placed_w, height=placed_h,
            )
            return

        if src_ratio > target_ratio:
            new_w = int(src_h * target_ratio)
            offset = (src_w - new_w) // 2
            box = (offset, 0, offset + new_w, src_h)
        else:
            new_h = int(src_w / target_ratio)
            offset = (src_h - new_h) // 2
            box = (0, offset, src_w, offset + new_h)

        cropped = im.crop(box)
        if cropped.mode not in ("RGB", "RGBA"):
            cropped = cropped.convert("RGBA")

        buf = io.BytesIO()
        cropped.save(buf, format="PNG")
        buf.seek(0)

    slide.shapes.add_picture(buf, left, top, width=width, height=height)


def _add_text(slide, left, top, width, height, text: str, *, font_name: str,
              font_size, color, bold: bool = False, align=PP_ALIGN.LEFT,
              vertical_anchor=MSO_ANCHOR.TOP) -> None:
    text = _clean_text(text)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = vertical_anchor
    # Render each line as its own paragraph so newlines are preserved
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color


def _add_solid_rect(slide, left, top, width, height, color) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def _draw_rainbow_stripe(slide, brand: Brand) -> None:
    """Brand-color stripe on content slides (Intro / Overview / Category).

    Dark themes: stripe sits at the bottom (against the dark bg it reads as a
    footer accent). Light themes: stripe sits under the title to break up the
    top white-space and match the Summary / AdditionalData chrome.
    """
    if brand.is_light:
        cfg_top = theme.LightTitle.rainbow_top
        seg_w = theme.LightTitle.rainbow_seg_w
        seg_h = theme.LightTitle.rainbow_h
        start_left = theme.LightTitle.rainbow_start_left
        count = theme.LightTitle.rainbow_count
    else:
        cfg_top = theme.Content.rainbow_top
        seg_w = theme.Content.rainbow_seg_w
        seg_h = theme.Content.rainbow_h
        start_left = theme.Content.rainbow_start_left
        count = theme.Content.rainbow_count

    for i, color in enumerate(brand.stripe_colors[:count]):
        left = start_left + seg_w * i
        _add_solid_rect(slide, left, cfg_top, seg_w, seg_h, color)


def _draw_content_title(slide, brand: Brand, title: str, *, color=None) -> None:
    _add_text(
        slide,
        theme.Content.title_left,
        theme.Content.title_top,
        theme.Content.title_w,
        theme.Content.title_h,
        title,
        font_name=brand.heading_font,
        font_size=theme.Content.title_font_pt,
        color=color or brand.text,
        bold=True,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def _draw_content_logo(slide, brand: Brand) -> None:
    if brand.logo_path and brand.logo_path.exists():
        _place_image(
            slide,
            brand.logo_path,
            theme.Content.logo_left,
            theme.Content.logo_top,
            theme.Content.logo_w,
            theme.Content.logo_h,
            fit=True,
        )


# ────────────────────────────────────────────────────────────────────────
# Cover
# ────────────────────────────────────────────────────────────────────────

def draw_cover(prs: Presentation, brand: Brand, *,
               title: str, subtitle: str, period: str,
               hero_override: Path | None = None) -> None:
    slide = _add_slide(prs, brand.background)

    _add_solid_rect(slide,
                    theme.Cover.stripe_left, theme.Cover.stripe_top,
                    theme.Cover.stripe_w, theme.Cover.stripe_h,
                    brand.accent)

    _add_text(slide,
              theme.Cover.title_left, theme.Cover.title_top,
              theme.Cover.title_w, theme.Cover.title_h,
              title,
              font_name=brand.heading_font,
              font_size=theme.Cover.title_font_pt,
              color=brand.text, bold=True)

    _add_text(slide,
              theme.Cover.period_left, theme.Cover.period_top,
              theme.Cover.period_w, theme.Cover.period_h,
              period,
              font_name=brand.heading_font,
              font_size=theme.Cover.period_font_pt,
              color=brand.text, bold=True)

    _add_text(slide,
              theme.Cover.subtitle_left, theme.Cover.subtitle_top,
              theme.Cover.subtitle_w, theme.Cover.subtitle_h,
              subtitle,
              font_name=brand.body_font,
              font_size=theme.Cover.subtitle_font_pt,
              color=brand.text)

    hero = brand.hero_path if (brand.hero_path and brand.hero_path.exists()) else hero_override
    if hero and Path(hero).exists():
        _place_image(slide, Path(hero),
                     theme.Cover.hero_left, theme.Cover.hero_top,
                     theme.Cover.hero_w, theme.Cover.hero_h)

    if brand.logo_path and brand.logo_path.exists():
        _place_image(slide, brand.logo_path,
                     theme.Cover.logo_left, theme.Cover.logo_top,
                     theme.Cover.logo_w, theme.Cover.logo_h,
                     fit=True)


# ────────────────────────────────────────────────────────────────────────
# Intro (per-account)
# ────────────────────────────────────────────────────────────────────────

def draw_intro(prs: Presentation, brand: Brand, *,
               title: str, body: str, previews: list[CategoryPreview]) -> None:
    """Per-account intro: title + body describing the categories + up to 4 images."""
    slide = _add_slide(prs, brand.background)
    _draw_content_title(slide, brand, title)
    _draw_content_logo(slide, brand)

    _add_text(slide,
              theme.Intro.body_left, theme.Intro.body_top,
              theme.Intro.body_w, theme.Intro.body_h,
              body,
              font_name=brand.body_font,
              font_size=theme.Intro.body_font_pt,
              color=brand.text,
              vertical_anchor=MSO_ANCHOR.MIDDLE)

    capped = previews[: theme.Intro.col_count]
    n = len(capped)
    for i, preview in enumerate(capped):
        _place_image(slide, preview.image_path,
                     theme.intro_col_left(i, n), theme.Intro.img_top,
                     theme.Intro.col_w, theme.Intro.img_h)

    _draw_rainbow_stripe(slide, brand)


# ────────────────────────────────────────────────────────────────────────
# Overview (5 metric circles)
# ────────────────────────────────────────────────────────────────────────

def _metric_font_size(number: str):
    n = len(number)
    if n <= 3:
        return Pt(40)
    if n <= 5:
        return Pt(28)
    return Pt(22)


def _add_metric_circle(slide, brand: Brand, left, top, diameter,
                       number: str, label: str,
                       sub_caption: str | None = None) -> None:
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    oval.fill.solid()
    oval.fill.fore_color.rgb = brand.primary
    oval.line.fill.background()

    tf = oval.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = number
    run.font.name = brand.heading_font
    run.font.size = _metric_font_size(number)
    run.font.bold = True
    # White inside the primary-color circle gives strong contrast on both dark
    # and light themes (circle is always brand.primary, not brand.background).
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    label_top_emu = top + diameter + theme.Overview.label_gap
    _add_text(slide, left, label_top_emu, diameter, theme.Overview.label_h,
              label,
              font_name=brand.body_font,
              font_size=theme.Overview.label_font_pt,
              color=brand.text,
              bold=True,
              align=PP_ALIGN.CENTER)

    if sub_caption:
        sub_top = label_top_emu + theme.Overview.label_h
        _add_text(slide, left, sub_top, diameter, theme.Overview.sub_caption_h,
                  sub_caption,
                  font_name=brand.body_font,
                  font_size=theme.Overview.sub_caption_font_pt,
                  color=brand.text,
                  align=PP_ALIGN.CENTER)


def draw_account_overview(prs: Presentation, brand: Brand, *,
                           title: str, metrics: list) -> None:
    """5 metric circles: Posts / Reels / Stories / Likes / Comments.

    Each metric is either (label, value) or (label, value, sub_caption).
    The Posts circle uses a sub-caption ("incl. N reels") to clarify that
    the value is the grand total of posts+reels.
    """
    slide = _add_slide(prs, brand.background)
    _draw_content_title(slide, brand, title)
    _draw_content_logo(slide, brand)

    for i, m in enumerate(metrics[: theme.Overview.circle_count]):
        if len(m) == 3:
            label, value, sub = m
        else:
            label, value = m
            sub = None
        _add_metric_circle(slide, brand,
                           theme.overview_circle_left(i),
                           theme.Overview.circle_top,
                           theme.Overview.circle_d,
                           value, label, sub_caption=sub)

    _draw_rainbow_stripe(slide, brand)


# ────────────────────────────────────────────────────────────────────────
# Category section (4 items per slide; auto-paginated)
# ────────────────────────────────────────────────────────────────────────

def draw_category_section(prs: Presentation, brand: Brand, *,
                          title: str, narrative: str,
                          items: list[CategoryItem]) -> None:
    """Paginate items and render the same narrative on every page.

    Kept for the visual-smoke spike (`scripts/_spike_report_render.py`) which
    exercises layouts with placeholder data and doesn't generate per-page
    narratives. Production reports go through the orchestrator
    (renderer.py), which calls `draw_category_page` directly with a
    per-page narrative.
    """
    per = theme.CategorySection.items_per_slide
    pages = [items[i:i + per] for i in range(0, len(items), per)] or [[]]
    total_pages = len(pages)
    for page_idx, page in enumerate(pages):
        draw_category_page(
            prs, brand,
            title=title, narrative=narrative, items=page,
            page_idx=page_idx, total_pages=total_pages,
        )


def draw_category_page(prs: Presentation, brand: Brand, *,
                       title: str, narrative: str,
                       items: list[CategoryItem],
                       page_idx: int = 0, total_pages: int = 1) -> None:
    """Render ONE slide for a single category page.

    The orchestrator handles pagination and supplies a per-page narrative
    (mentioning only the items on this page). When `total_pages > 1`, the
    title gets a "(i/N)" suffix. Items are horizontally centered when fewer
    than `items_per_slide` are passed.
    """
    slide = _add_slide(prs, brand.background)
    page_title = (
        title if total_pages == 1
        else f"{title} ({page_idx + 1}/{total_pages})"
    )
    _draw_content_title(slide, brand, page_title)
    _draw_content_logo(slide, brand)

    # Per-page prose — vertically centered between title and images
    _add_text(slide,
              theme.CategorySection.narrative_left,
              theme.CategorySection.narrative_top,
              theme.CategorySection.narrative_w,
              theme.CategorySection.narrative_h,
              narrative,
              font_name=brand.body_font,
              font_size=theme.CategorySection.narrative_font_pt,
              color=brand.text,
              vertical_anchor=MSO_ANCHOR.MIDDLE)

    n = len(items)
    for col, item in enumerate(items):
        _place_image(slide, item.image_path,
                     theme.category_col_left(col, n),
                     theme.CategorySection.img_top,
                     theme.CategorySection.col_w,
                     theme.CategorySection.img_h)
        _add_text(slide,
                  theme.category_col_left(col, n),
                  theme.CategorySection.caption_top,
                  theme.CategorySection.col_w,
                  theme.CategorySection.caption_h,
                  _split_caption(item.narrative),
                  font_name=brand.body_font,
                  font_size=theme.CategorySection.caption_font_pt,
                  color=brand.text,
                  align=PP_ALIGN.CENTER,
                  vertical_anchor=MSO_ANCHOR.TOP)

    _draw_rainbow_stripe(slide, brand)


# ────────────────────────────────────────────────────────────────────────
# Light-theme helpers
# ────────────────────────────────────────────────────────────────────────

def _draw_light_header(slide, brand: Brand, title: str) -> None:
    """Title + logo + rainbow stripe — shared by Summary + AdditionalData.

    Uses the same 8-segment rainbow pattern as the dark-theme bottom rainbow,
    mirrored under the heading to maintain design language across themes.
    """
    _add_text(slide,
              theme.LightTitle.title_left, theme.LightTitle.title_top,
              theme.LightTitle.title_w, theme.LightTitle.title_h,
              title,
              font_name=brand.heading_font,
              font_size=theme.LightTitle.title_font_pt,
              color=brand.secondary,
              bold=True,
              vertical_anchor=MSO_ANCHOR.MIDDLE)

    # Rainbow stripe under title (same geometry as bottom rainbow on dark slides)
    for i, color in enumerate(brand.stripe_colors[: theme.LightTitle.rainbow_count]):
        left = theme.LightTitle.rainbow_start_left + theme.LightTitle.rainbow_seg_w * i
        _add_solid_rect(slide, left,
                        theme.LightTitle.rainbow_top,
                        theme.LightTitle.rainbow_seg_w,
                        theme.LightTitle.rainbow_h,
                        color)

    if brand.logo_path and brand.logo_path.exists():
        _place_image(slide, brand.logo_path,
                     theme.LightTitle.logo_left, theme.LightTitle.logo_top,
                     theme.LightTitle.logo_w, theme.LightTitle.logo_h,
                     fit=True)


def _add_metric_card(slide, brand: Brand, left, top, width, height,
                     card: MetricCard, value_font_pt) -> None:
    """One card on Summary / AdditionalData. Light theme: navy text on white."""
    # Subheading label (top of card)
    label_h = Inches(0.40)
    _add_text(slide, left, top, width, label_h,
              card.label,
              font_name=brand.heading_font,
              font_size=theme.Summary.label_font_pt,
              color=brand.primary,     # accent color for subheadings (orange-ish)
              bold=True)

    # Big primary value (optionally followed by a smaller, non-bold aux run
    # in the same paragraph so trailing detail like "(10 reels)" doesn't wrap
    # the bold value onto a second line and clobber the caption underneath).
    value_top = top + label_h
    value_h = Inches(0.95)
    if card.value_aux:
        box = slide.shapes.add_textbox(left, value_top, width, value_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r_main = p.add_run()
        r_main.text = _clean_text(card.value)
        r_main.font.name = brand.heading_font
        r_main.font.size = value_font_pt
        r_main.font.bold = True
        r_main.font.color.rgb = brand.secondary
        r_aux = p.add_run()
        r_aux.text = " " + _clean_text(card.value_aux)
        r_aux.font.name = brand.body_font
        r_aux.font.size = Pt(16)
        r_aux.font.bold = False
        r_aux.font.color.rgb = brand.secondary
    else:
        _add_text(slide, left, value_top, width, value_h,
                  card.value,
                  font_name=brand.heading_font,
                  font_size=value_font_pt,
                  color=brand.secondary,   # navy
                  bold=True,
                  vertical_anchor=MSO_ANCHOR.TOP)

    # Caption (smaller detail) — fills the rest of the card
    caption_top = value_top + value_h
    caption_h = int(height) - int(label_h) - int(value_h)
    if card.caption and caption_h > 0:
        _add_text(slide, left, caption_top, width, Emu(caption_h),
                  card.caption,
                  font_name=brand.body_font,
                  font_size=theme.Summary.caption_font_pt,
                  color=brand.secondary)


def draw_account_summary(prs: Presentation, brand: Brand, *,
                          title: str, cards: list[MetricCard]) -> None:
    """Light-theme data digest: 2x3 grid of subheading + value + caption cards."""
    slide = _add_slide(prs, LIGHT_BG)
    _draw_light_header(slide, brand, title)

    for i, card in enumerate(cards[: theme.Summary.grid_rows * theme.Summary.grid_cols]):
        row, col = divmod(i, theme.Summary.grid_cols)
        left, top, w, h = theme.summary_card_pos(row, col)
        _add_metric_card(slide, brand, left, top, w, h, card,
                         value_font_pt=theme.Summary.value_font_pt)


def draw_additional_data(prs: Presentation, brand: Brand, *,
                          title: str, cards: list[MetricCard]) -> None:
    """Light-theme extended-data slide. Same shape as Summary; placeholder values OK."""
    slide = _add_slide(prs, LIGHT_BG)
    _draw_light_header(slide, brand, title)

    for i, card in enumerate(cards[: theme.AdditionalData.grid_rows * theme.AdditionalData.grid_cols]):
        row, col = divmod(i, theme.AdditionalData.grid_cols)
        left, top, w, h = theme.summary_card_pos(row, col)  # same grid math
        _add_metric_card(slide, brand, left, top, w, h, card,
                         value_font_pt=theme.AdditionalData.value_font_pt)
