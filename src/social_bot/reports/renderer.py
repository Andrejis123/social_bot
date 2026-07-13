"""Top-level report orchestration.

Single entry point: `generate_report(client_slug, period)` runs
`data.load_report_data` → per-(account, category) `synthesis.synthesize_category`
→ assembles the slide sequence and saves a .pptx.

Slide sequence (locked 2026-05-26):

    Cover (once, client-level)
    for each active account in client.yaml order:
        Intro            (categories description + 4 previews)
        Overview         (5 metric circles)
        Category × N     (paginated 4 items per slide; ordered by post count desc)
        Summary          (6 cards — light theme)
        Additional Data  (6 placeholder cards — light theme; see project_additional_data_future.md)

Decisions wired here (chosen with user 2026-05-28):
- Intro body: deterministic template "N posts across A (n), B (n), ..."  (no LLM call)
- Summary "Collaborations" card: synthesis-derived (cluster count + comma-joined titles)
- Additional Data: static placeholder strings (slide layout exercised; data wiring deferred)
- (Uncategorized) bucket: skipped entirely (see project_uncategorized_policy.md)
- Output: local .pptx only for now; Supabase upload + email follow-up deferred

Image fallback for best_post_id:
- synthesis.py picks a best_post_id but doesn't know about image availability.
- If `posts_by_id[best_post_id].hero_image_path is None`, this module re-picks
  the highest-engagement post in the same cluster that has an image.
- If the entire cluster has no images (rare: reel-only cluster with expired
  IG signed URLs for all covers), the cluster is dropped from the slide.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING, cast

from pptx import Presentation

from .. import drive
from ..config import get_settings
from ..db import queries
from ..logging import get_logger
from ..notifications import telegram
from ..storage.reports import UploadedReport, upload_report
from ..storage.synthesis import load_latest_synthesis_artifact, save_synthesis_artifact
from . import theme
from .brand import Brand
from .data import (
    UNCATEGORIZED,
    AccountData,
    Period,
    PostRow,
    ReportData,
    load_report_data,
    reel_term,
)
from .layouts import (
    CategoryItem,
    CategoryPreview,
    MetricCard,
    draw_account_overview,
    draw_account_summary,
    draw_additional_data,
    draw_category_page,
    draw_cover,
    draw_intro,
    format_metric,
)
from .synthesis import (
    PROMPT_VERSIONS,
    CategorySynthesis,
    ClusterItem,
    synthesize_category,
    synthesize_page_narrative,
)

if TYPE_CHECKING:
    # `pptx.Presentation` (imported above) is the factory function used at
    # runtime; the actual class for annotations lives in pptx.presentation.
    from pptx.presentation import Presentation as PresentationDoc

log = get_logger(__name__)

DEFAULT_OUT_DIR = Path("/tmp/reports")


# ─────────────────────────────────────────────────────────────────────
# Public entry point
# ─────────────────────────────────────────────────────────────────────

@dataclass(slots=True)
class _BuiltReport:
    path: Path
    report: ReportData
    slide_count: int


PLATFORM_LABELS = {
    "instagram": "Instagram",
    "facebook": "Facebook",
    "tiktok": "TikTok",
}


def _build_report(
    client_slug: str,
    period: Period,
    *,
    out_dir: Path = DEFAULT_OUT_DIR,
    platform: str | None = None,
    reuse_synthesis: bool = False,
) -> _BuiltReport:
    """Render the deck and return path + ReportData + slide count.

    Internal helper so `generate_report` and `publish_report` share one
    Supabase fetch + one Presentation instance. When `platform` is set, only
    that platform's accounts are included (a standalone single-platform deck).

    When `reuse_synthesis=True`, skips all LLM calls and loads the most recent
    synthesis artifact from Supabase for this client+period+platform. Use this
    to render color/layout variants without burning new LLM passes.
    """
    platform_key = platform or "instagram"
    settings = get_settings()
    out_dir.mkdir(parents=True, exist_ok=True)

    report = load_report_data(client_slug, period, platform=platform)
    brand = Brand.load(client_slug)

    # Load precomputed synthesis when requested.
    precomputed_all: dict[str, dict[str, CategorySynthesis]] | None = None
    if reuse_synthesis:
        blob = load_latest_synthesis_artifact(
            client_slug=client_slug,
            period_label=period.label,
            platform=platform_key,
        )
        if blob is None:
            raise RuntimeError(
                f"No synthesis artifact found for {client_slug} / {period.label} / {platform_key}. "
                "Run without --reuse-synthesis first to generate one."
            )
        precomputed_all = {
            handle: {
                cat: CategorySynthesis.from_dict(synth_dict)
                for cat, synth_dict in cats.items()
            }
            for handle, cats in blob.items()
        }
        log.info(
            "report.reuse_synthesis",
            client=client_slug, period=period.label, platform=platform_key,
            accounts=len(precomputed_all),
        )

    prs = Presentation()
    prs.slide_width = theme.SLIDE_W
    prs.slide_height = theme.SLIDE_H

    # Cover — once. Fall back to highest-engagement healed post image across
    # the report when brand has no static hero_path (avoids a blank right side).
    hero_override: Path | None = None
    if not (brand.hero_path and brand.hero_path.exists()):
        all_posts = [
            p for acct in report.accounts
            for posts in acct.posts_by_category.values()
            for p in posts
            if p.hero_image_path
        ]
        if all_posts:
            top = max(all_posts,
                      key=lambda p: p.engagement)
            hero_override = top.hero_image_path

    platform_label = PLATFORM_LABELS.get(platform_key)
    cover_title = (
        f"{platform_label} Activity Monitoring" if platform_label
        else "Social Media Monitoring"
    )
    draw_cover(
        prs, brand,
        title=cover_title,
        subtitle=f"for {report.client_name}",
        period=period.label,
        hero_override=hero_override,
    )

    # Per-account block. Collect synthesis when running fresh (to persist as artifact).
    collected_synth: dict[str, dict[str, CategorySynthesis]] = {}
    for account in report.accounts:
        precomputed = precomputed_all.get(account.handle) if precomputed_all else None
        synth_by_cat = _render_account(
            prs, brand, account, report, period, precomputed=precomputed,
        )
        if precomputed_all is None:
            collected_synth[account.handle] = synth_by_cat

    if precomputed_all is None:
        try:
            save_synthesis_artifact(
                client_slug=client_slug,
                period_label=period.label,
                platform=platform_key,
                model=settings.gemini_model,
                prompt_versions=PROMPT_VERSIONS,
                artifact={
                    handle: {cat: s.to_dict() for cat, s in cats.items()}
                    for handle, cats in collected_synth.items()
                },
            )
        except Exception as exc:
            log.warning("report.synthesis_artifact_save_failed", error=str(exc))

    # Filename doubles as the Supabase Storage key — strip any character that
    # Supabase's `InvalidKey` validator rejects (en/em dashes, etc.).
    safe_period = (
        period.label
        .replace(" ", "_")
        .replace("/", "-")
        .replace("–", "-")   # en-dash from period labels
        .replace("—", "-")   # em-dash, defense in depth
    )
    platform_tag = f"{platform}_" if platform else ""
    out_path = out_dir / f"{client_slug}_{platform_tag}{safe_period}.pptx"
    prs.save(str(out_path))
    slide_count = len(prs.slides)
    log.info(
        "report.saved", client=client_slug, period=period.label,
        path=str(out_path), accounts=len(report.accounts),
        slides=slide_count,
    )
    return _BuiltReport(path=out_path, report=report, slide_count=slide_count)


def generate_report(
    client_slug: str,
    period: Period,
    *,
    out_dir: Path = DEFAULT_OUT_DIR,
    platform: str | None = None,
    reuse_synthesis: bool = False,
) -> Path:
    """Build the per-client monthly deck. Returns the saved .pptx path.

    When `platform` is set (e.g. "facebook"), only that platform's accounts
    are included — a standalone single-platform deck.
    """
    return _build_report(
        client_slug, period, out_dir=out_dir, platform=platform,
        reuse_synthesis=reuse_synthesis,
    ).path


def publish_report(
    client_slug: str,
    period: Period,
    *,
    out_dir: Path = DEFAULT_OUT_DIR,
    platform: str | None = None,
    reuse_synthesis: bool = False,
) -> tuple[Path, UploadedReport]:
    """Render, upload to Supabase, and notify on Telegram.

    Notification failures are swallowed by telegram.send (per its contract);
    upload failures bubble.
    """
    built = _build_report(
        client_slug, period, out_dir=out_dir, platform=platform,
        reuse_synthesis=reuse_synthesis,
    )
    uploaded = upload_report(client_slug, built.path)

    # Record the successful report (best-effort): this row is what the
    # report-gated archive cron checks before archiving the period. A
    # recording failure must not break an otherwise-delivered report; the gap
    # surfaces later as the archive gate skipping that client + alerting.
    try:
        queries.record_report_run(
            client_slug=client_slug,
            period_start=period.start.date(),
            period_end=period.end.date(),
            platform=platform,
            slide_count=built.slide_count,
            bytes_size=uploaded.bytes_size,
        )
    except Exception as exc:
        log.warning("report.record_run_failed", client=client_slug, error=str(exc))

    # Drive upload is best-effort: failure must not break the existing
    # Supabase + Telegram path (which is what cron relies on for delivery).
    try:
        drive_result = drive.upload_report(client_slug, built.path)
        log.info(
            "report.drive_uploaded",
            client=client_slug,
            file_id=drive_result["id"],
            url=drive_result["webViewLink"],
        )
    except Exception as exc:
        log.warning("report.drive_upload_failed", client=client_slug, error=str(exc))

    telegram.notify_report_generated(
        client_name=built.report.client_name,
        period_label=period.label,
        slide_count=built.slide_count,
        bytes_size=uploaded.bytes_size,
        signed_url=uploaded.signed_url,
    )

    return built.path, uploaded


# ─────────────────────────────────────────────────────────────────────
# Per-account orchestration
# ─────────────────────────────────────────────────────────────────────

def _render_account(
    prs: PresentationDoc, brand: Brand,
    account: AccountData, report: ReportData, period: Period,
    *,
    precomputed: dict[str, CategorySynthesis] | None = None,
) -> dict[str, CategorySynthesis]:
    handle = account.handle
    display = f"@{handle}"

    # Drop the Uncategorized bucket from anything report-facing (see
    # project_uncategorized_policy.md).
    real_cats: list[tuple[str, list[PostRow]]] = [
        (cat, posts) for cat, posts in account.posts_by_category.items()
        if cat != UNCATEGORIZED
    ]

    # ── Intro
    draw_intro(
        prs, brand,
        title=f"{display}: Intro",
        body=_build_intro_body(account, real_cats),
        previews=_build_intro_previews(account, real_cats),
    )

    # ── Overview
    draw_account_overview(
        prs, brand,
        title=f"{display}: Activity Overview",
        metrics=_build_overview_metrics(account),
    )

    # ── Categories: synthesize each, render in count-desc order (data.py
    # already orders them). Capture synthesis for Summary's collab card.
    # used_paths is shared across categories so image picks dedupe globally
    # (matters when heal failures leave only a handful of healed images).
    synth_by_cat: dict[str, CategorySynthesis] = {}
    account_posts = [p for _, posts in real_cats for p in posts]
    used_paths: set[str] = set()
    for cat, posts in real_cats:
        if precomputed and cat in precomputed:
            synth = precomputed[cat]
        else:
            synth = synthesize_category(
                client_slug=report.client_slug,
                period_label=period.label,
                brand_label=display,   # "@handle" — Pass 0 v2 references this verbatim
                category=cat,
                posts=posts,
            )
        synth_by_cat[cat] = synth

        items = _items_from_synthesis(synth, posts, account_posts, used_paths)
        if not items:
            log.warning(
                "report.category_no_items",
                client=report.client_slug, account=handle, category=cat,
            )
            continue

        # Paginate after image resolution so per-page narratives match the
        # items that actually survived. Single-page categories keep the
        # existing Pass-0 category narrative (zero extra LLM calls); multi-
        # page categories get a fresh per-page Pass-0 narrative against the
        # subset of posts whose clusters land on that page.
        per = theme.CategorySection.items_per_slide
        pages = [items[i:i + per] for i in range(0, len(items), per)]
        total_pages = len(pages)
        posts_by_id = {p.id: p for p in posts}

        for page_idx, page_items in enumerate(pages):
            if total_pages == 1:
                page_narrative = synth.category_narrative
            else:
                page_posts: list[PostRow] = []
                seen: set[str] = set()
                for it in page_items:
                    for pid in it.post_ids:
                        if pid in seen:
                            continue
                        post = posts_by_id.get(pid)
                        if post is None:
                            continue
                        page_posts.append(post)
                        seen.add(pid)
                page_narrative = synthesize_page_narrative(
                    client_slug=report.client_slug,
                    period_label=period.label,
                    brand_label=display,   # "@handle"
                    category=cat,
                    posts=page_posts,
                ) or synth.category_narrative

            draw_category_page(
                prs, brand,
                title=f"{display}: {cat}",
                narrative=page_narrative,
                items=page_items,
                page_idx=page_idx,
                total_pages=total_pages,
            )

    # ── Summary (light theme)
    draw_account_summary(
        prs, brand,
        title=f"{display}: Summary",
        cards=_build_summary_cards(account, real_cats, synth_by_cat, period),
    )

    # ── Additional Data (light theme, static placeholders for v1)
    draw_additional_data(
        prs, brand,
        title=f"{display}: Additional Data",
        cards=_additional_data_placeholders(),
    )

    return synth_by_cat


# ─────────────────────────────────────────────────────────────────────
# Intro body — deterministic template
# ─────────────────────────────────────────────────────────────────────

def _build_intro_previews(
    account: AccountData,
    real_cats: list[tuple[str, list[PostRow]]],
) -> list[CategoryPreview]:
    """Pick 3-4 images for the Intro strip.

    Priority: one image per top category (data.py already does this and stores
    them in `intro_previews`, ordered by category size desc, Uncategorized
    excluded). If fewer than `min_previews` distinct-category previews are
    available (low-activity accounts), top up with extra high-engagement posts
    from the largest category — duplicates are acceptable since the goal is
    visual fill, not strict uniqueness across the deck.
    """
    base = [
        CategoryPreview(name=p.name, image_path=p.image_path)
        for p in account.intro_previews if p.name != UNCATEGORIZED
    ]

    # Track which post images are already on the slide to avoid placing the
    # same image twice when filling.
    used_paths: set[str] = {str(p.image_path) for p in base}

    if len(base) < theme.Intro.min_previews:
        # Walk posts from the most-populated category first; for each, take
        # any post with a hero image not already shown.
        filler_pool: list[PostRow] = []
        for _, posts in real_cats:
            filler_pool.extend(
                sorted(
                    [p for p in posts if p.hero_image_path is not None],
                    key=lambda p: p.engagement,
                    reverse=True,
                )
            )
        for post in filler_pool:
            if len(base) >= theme.Intro.min_previews:
                break
            path_str = str(post.hero_image_path)
            if path_str in used_paths:
                continue
            base.append(CategoryPreview(
                name=post.ai_category or "",
                # filler_pool only holds posts with a non-None hero image.
                image_path=cast(Path, post.hero_image_path),
            ))
            used_paths.add(path_str)

        # Last-resort fill: if still under min_previews (heal-failure cascade
        # left too few unique images), allow image duplicates so the strip
        # reaches the spec-locked minimum count.
        if len(base) < theme.Intro.min_previews and filler_pool:
            for post in filler_pool:
                if len(base) >= theme.Intro.min_previews:
                    break
                base.append(CategoryPreview(
                    name=post.ai_category or "",
                    image_path=cast(Path, post.hero_image_path),
                ))

    return base[: theme.Intro.col_count]


def _build_intro_body(
    account: AccountData,
    real_cats: list[tuple[str, list[PostRow]]],
) -> str:
    """Compose the Intro body from category names + counts.

    Examples:
      multi  → "35 posts this period across categories: Events (17), Competition (8), Ongoing (8), Collaborations (2)."
      single → "5 posts this period across category: Ongoing (5)."
    """
    total = sum(len(posts) for _, posts in real_cats)
    if total == 0:
        return f"No posts for this period from @{account.handle}."

    parts = [f"{cat} ({len(posts)})" for cat, posts in real_cats]
    cats_word = "category" if len(parts) == 1 else "categories"
    plural_p = "posts" if total != 1 else "post"
    return f"{total} {plural_p} this period across {cats_word}: {', '.join(parts)}."


# ─────────────────────────────────────────────────────────────────────
# Overview metrics
# ─────────────────────────────────────────────────────────────────────

def _build_overview_metrics(account: AccountData) -> list[tuple]:
    """5 circles: Total posts including reels / Reels / Stories / Likes / Comments
    (reels-terminology per platform, see data.reel_term).

    First circle label spells out the relationship explicitly so the user
    isn't left guessing whether reels are double-counted (they are).
    """
    term = reel_term(account.platform)
    total_posts = account.total_posts + account.total_reels
    return [
        (f"Total posts including {term}", format_metric(total_posts)),
        (term.capitalize(), format_metric(account.total_reels)),
        ("Stories", format_metric(account.total_stories)),
        ("Total Likes", format_metric(account.total_likes)),
        ("Total Comments", format_metric(account.total_comments)),
    ]


# ─────────────────────────────────────────────────────────────────────
# Cluster → CategoryItem (with image fallback)
# ─────────────────────────────────────────────────────────────────────

def _items_from_synthesis(
    synth: CategorySynthesis,
    posts: list[PostRow],
    account_posts: list[PostRow],
    used_paths: set[str],
) -> list[CategoryItem]:
    posts_by_id = {p.id: p for p in posts}
    items: list[CategoryItem] = []
    for cluster in synth.items:
        chosen = _resolve_cluster_image(
            cluster, posts_by_id, account_posts, used_paths,
        )
        if chosen is None:
            log.warning(
                "report.cluster_dropped_no_image",
                cluster=cluster.title, best_post_id=cluster.best_post_id[:8],
            )
            continue
        used_paths.add(str(chosen))
        items.append(CategoryItem(
            title=cluster.title,
            narrative=cluster.narrative,
            image_path=chosen,
            post_ids=list(cluster.post_ids),
        ))
    return items


def _resolve_cluster_image(
    cluster: ClusterItem,
    posts_by_id: dict[str, PostRow],
    account_posts: list[PostRow],
    used_paths: set[str],
) -> Path | None:
    """Pick an image path for a cluster with layered fallbacks.

    Order: (1) synthesis best_post if image-bearing and unused; (2) highest-
    engagement category post with unused image; (3) highest-engagement account
    post with unused image (covers heal-failure cascades where a whole bucket
    has no images); (4) any category image (allow duplicate); (5) any account
    image (allow duplicate). Returns None only when the entire account has
    zero healed images.
    """
    best = posts_by_id.get(cluster.best_post_id)
    if best and best.hero_image_path and str(best.hero_image_path) not in used_paths:
        return best.hero_image_path

    cat_pool = sorted(
        [p for p in posts_by_id.values() if p.hero_image_path],
        key=lambda p: p.engagement,
        reverse=True,
    )
    for p in cat_pool:
        if str(p.hero_image_path) not in used_paths:
            return p.hero_image_path

    acct_pool = sorted(
        [p for p in account_posts if p.hero_image_path],
        key=lambda p: p.engagement,
        reverse=True,
    )
    for p in acct_pool:
        if str(p.hero_image_path) not in used_paths:
            return p.hero_image_path

    # All unique slots exhausted — allow duplicate from category, then account.
    if cat_pool:
        return cat_pool[0].hero_image_path
    if acct_pool:
        return acct_pool[0].hero_image_path
    return None


# ─────────────────────────────────────────────────────────────────────
# Summary cards
# ─────────────────────────────────────────────────────────────────────

def _build_summary_cards(
    account: AccountData,
    real_cats: list[tuple[str, list[PostRow]]],
    synth_by_cat: dict[str, CategorySynthesis],
    period: Period,
) -> list[MetricCard]:
    """Six cards in the locked order:
    Volume / Collaborations / Post cadence / Story cadence / Likes / Comments.
    """
    total_posts = account.total_posts + account.total_reels
    weeks = max(1, _weeks_in_period(period))
    days = max(1, _days_in_period(period))

    # Volume
    if real_cats:
        top_cat, top_posts = real_cats[0]
        if len(real_cats) > 1:
            second_cat, second_posts = real_cats[1]
            volume_caption = (
                f"Led by {top_cat} ({len(top_posts)}) "
                f"and {second_cat} ({len(second_posts)})."
            )
        else:
            volume_caption = f"All {top_cat} this period."
    else:
        volume_caption = "No activity in the period."

    reels_aux = (
        f"including {account.total_reels} {reel_term(account.platform)}"
        if account.total_reels
        else ""
    )
    volume = MetricCard(
        label="Volume",
        value=f"{total_posts} posts",
        caption=volume_caption,
        value_aux=reels_aux,
    )

    # Collaborations — pulled from synthesis if a "Collaborations" category exists.
    # Cluster titles can themselves contain commas (e.g. "Art Space by Ploom,
    # Magdalena Šťastníková"), so use " / " between cluster titles to keep
    # cluster boundaries readable.
    collab_synth = _find_collab_synthesis(synth_by_cat)
    if collab_synth and collab_synth.items:
        titles = [c.title for c in collab_synth.items]
        collab = MetricCard(
            label="Collaborations",
            value=f"{len(titles)} distinct",
            caption=" / ".join(titles) + ".",
        )
    else:
        collab = MetricCard(
            label="Collaborations",
            value="None tracked",
            caption="No partnerships surfaced during this period.",
        )

    # Cadence
    post_cadence = MetricCard(
        label="Post cadence",
        value=f"{total_posts / weeks:.1f} / week",
        caption=f"Across {weeks} weeks in the period.",
    )
    if account.total_stories:
        story_cadence = MetricCard(
            label="Story cadence",
            value=f"{account.total_stories / days:.1f} / day",
            caption=f"{account.total_stories} {'story' if account.total_stories == 1 else 'stories'} across {days} days.",
        )
    else:
        story_cadence = MetricCard(
            label="Story cadence",
            value="No stories",
            caption="Stories pipeline paused or account inactive.",
        )

    # Likes
    avg_likes = account.total_likes // total_posts if total_posts else 0
    likes = MetricCard(
        label="Likes",
        value=f"{format_metric(avg_likes)} avg",
        caption=f"{format_metric(account.total_likes)} total across {total_posts} posts.",
    )

    # Comments
    avg_comments = account.total_comments // total_posts if total_posts else 0
    comments = MetricCard(
        label="Comments",
        value=f"{format_metric(avg_comments)} avg",
        caption=f"{format_metric(account.total_comments)} total across {total_posts} posts.",
    )

    return [volume, collab, post_cadence, story_cadence, likes, comments]


def _find_collab_synthesis(synth_by_cat: dict[str, CategorySynthesis]) -> CategorySynthesis | None:
    """Case-insensitive lookup for a 'Collaborations' (or similar) category.

    Different clients name this differently (ecig: 'Collaborations'; agape may
    not have one at all). Tries a few common spellings.
    """
    candidates = ("collaborations", "collaboration", "partnerships", "partners",
                  "influencer", "influencers")
    for cat_name, synth in synth_by_cat.items():
        if cat_name.lower() in candidates:
            return synth
    return None


def _additional_data_placeholders() -> list[MetricCard]:
    """Six static cards exercising the AdditionalData layout. Real values
    require either velocity-scrape infra (see project_additional_data_future.md)
    or have been chosen as deferred-tier offerings."""
    return [
        MetricCard("Posting time window", "TBD",
                   "Time-of-day distribution from posts.posted_at."),
        MetricCard("Story time window", "TBD",
                   "Time-of-day distribution from stories.posted_at."),
        MetricCard("Top day of week", "TBD",
                   "Day with the highest average engagement."),
        MetricCard("Fastest-growing post", "TBD",
                   "Needs ≥2 metric snapshots within first 24h."),
        MetricCard("Fastest-commenting post", "TBD",
                   "Needs ≥2 metric snapshots within first 24h."),
        MetricCard("Hashtag diversity", "TBD",
                   "Unique hashtags across the period."),
    ]


# ─────────────────────────────────────────────────────────────────────
# Period helpers
# ─────────────────────────────────────────────────────────────────────

def _days_in_period(period: Period) -> int:
    return max(1, (period.end.date() - period.start.date()).days + 1)


def _weeks_in_period(period: Period) -> int:
    return max(1, _days_in_period(period) // 7)
